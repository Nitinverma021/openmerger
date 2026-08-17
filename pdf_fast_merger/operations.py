"""Standalone, non-destructive PDF page operations used alongside merging."""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import os
import shutil
import subprocess
import sys
import tempfile
from collections.abc import Iterable
from pathlib import Path
from typing import Literal

ImagePagePreset = Literal["original", "a4", "letter", "poster", "social", "banner"]
ImageFit = Literal["contain", "cover"]

PAGE_SIZES_INCHES: dict[ImagePagePreset, tuple[float, float] | None] = {
    "original": None,
    "a4": (8.2677, 11.6929),
    "letter": (8.5, 11.0),
    "poster": (11.0, 17.0),
    "social": (8.0, 10.0),
    "banner": (13.3333, 7.5),
}


def _ensure_new_output(source_path: Path, output_path: Path) -> None:
    if source_path.expanduser().resolve() == output_path.expanduser().resolve():
        raise ValueError("Choose a different output file; source PDFs are never overwritten.")


def read_image_metadata(image_path: Path) -> dict[str, object]:
    """Return a readable, local-only metadata report for an image."""
    from PIL import ExifTags, Image

    with Image.open(image_path) as image:
        report: dict[str, object] = {
            "File": image_path.name,
            "Format": image.format or image_path.suffix.lstrip(".").upper(),
            "File size": f"{image_path.stat().st_size:,} bytes",
            "Dimensions": f"{image.width} × {image.height} px",
            "Color mode": image.mode,
        }
        if image.info.get("dpi"):
            report["DPI"] = " × ".join(str(round(float(value), 2)) for value in image.info["dpi"])
        exif = image.getexif()
        for key, value in exif.items():
            name = ExifTags.TAGS.get(key, f"EXIF {key}")
            if name == "GPSInfo":
                try:
                    gps = exif.get_ifd(ExifTags.IFD.GPSInfo)
                except (AttributeError, KeyError, TypeError):
                    gps = value if isinstance(value, dict) else {}
                report["GPS metadata"] = {ExifTags.GPSTAGS.get(item, str(item)): _metadata_value(entry) for item, entry in gps.items()}
            else:
                report[name] = _metadata_value(value)
        xmp = image.info.get("xmp")
        if xmp:
            report["XMP metadata"] = f"Present ({len(xmp)} bytes)"
    return report


def _metadata_value(value: object) -> object:
    if isinstance(value, bytes):
        return f"Binary data ({len(value)} bytes)"
    if isinstance(value, dict):
        return {str(key): _metadata_value(item) for key, item in value.items()}
    if isinstance(value, (list, tuple)):
        return [_metadata_value(item) for item in value]
    return str(value)


def image_metadata_report(image_path: Path) -> str:
    return json.dumps(read_image_metadata(image_path), indent=2, ensure_ascii=False)


def image_metadata_reports(image_paths: Iterable[Path]) -> str:
    """Return one JSON document describing every selected image."""
    reports = [read_image_metadata(path) for path in image_paths]
    if not reports:
        raise ValueError("Choose at least one image.")
    return json.dumps({"images": reports}, indent=2, ensure_ascii=False)


def remove_image_metadata(source_path: Path, output_path: Path) -> None:
    """Create a visually equivalent image copy without EXIF/XMP/private metadata."""
    from PIL import Image

    _ensure_new_output(source_path, output_path)
    with Image.open(source_path) as source:
        clean = source.copy()
        try:
            suffix = output_path.suffix.casefold()
            save_options: dict[str, object] = {}
            if suffix in {".jpg", ".jpeg"}:
                if clean.mode not in {"RGB", "L"}:
                    converted = clean.convert("RGB")
                    clean.close()
                    clean = converted
                save_options = {"quality": 95, "subsampling": 0, "optimize": True}
            partial = output_path.with_suffix(output_path.suffix + ".partial")
            formats = {".jpg": "JPEG", ".jpeg": "JPEG", ".png": "PNG", ".webp": "WEBP", ".bmp": "BMP", ".tif": "TIFF", ".tiff": "TIFF"}
            image_format = formats.get(suffix)
            if not image_format:
                raise ValueError("Output must use JPG, PNG, WebP, BMP, or TIFF.")
            clean.save(partial, format=image_format, **save_options)
            os.replace(partial, output_path)
        finally:
            clean.close()


def convert_images(
    source_paths: Iterable[Path],
    output_folder: Path,
    output_format: Literal["jpg", "png", "webp", "bmp", "tiff"] = "png",
    max_width: int = 0,
    max_height: int = 0,
    quality: int = 90,
) -> list[Path]:
    """Convert, optionally resize, and locally optimize one or more images."""
    from PIL import Image, ImageOps

    _register_heif_support()
    paths = list(source_paths)
    if not paths:
        raise ValueError("Choose at least one image.")
    if max_width < 0 or max_height < 0 or not 1 <= quality <= 100:
        raise ValueError("Image dimensions must be positive and quality must be between 1 and 100.")
    formats = {"jpg": ("JPEG", ".jpg"), "png": ("PNG", ".png"), "webp": ("WEBP", ".webp"), "bmp": ("BMP", ".bmp"), "tiff": ("TIFF", ".tiff")}
    image_format, suffix = formats[output_format]
    output_folder.mkdir(parents=True, exist_ok=True)
    outputs: list[Path] = []
    for source_path in paths:
        target = output_folder / f"{source_path.stem}_converted{suffix}"
        if target.expanduser().resolve() == source_path.expanduser().resolve():
            target = output_folder / f"{source_path.stem}_converted_copy{suffix}"
        with Image.open(source_path) as source:
            image = ImageOps.exif_transpose(source)
            try:
                if max_width or max_height:
                    bounds = (max_width or image.width, max_height or image.height)
                    image.thumbnail(bounds, Image.Resampling.LANCZOS)
                if image_format == "JPEG" and image.mode not in {"RGB", "L"}:
                    converted = image.convert("RGB")
                    image.close()
                    image = converted
                partial = target.with_suffix(target.suffix + ".partial")
                options: dict[str, object] = {"quality": quality} if image_format in {"JPEG", "WEBP"} else {}
                image.save(partial, format=image_format, **options)
                os.replace(partial, target)
                outputs.append(target)
            finally:
                image.close()
    return outputs


def remove_image_background(source_path: Path, output_path: Path) -> None:
    """Create a transparent PNG by separating a centered foreground from its background locally."""
    import cv2
    import numpy as np

    _ensure_new_output(source_path, output_path)
    if output_path.suffix.casefold() != ".png":
        raise ValueError("Background removal creates a transparent PNG; choose a .png output name.")
    image = cv2.imread(str(source_path), cv2.IMREAD_UNCHANGED)
    if image is None:
        raise ValueError("Unable to read this image.")
    if image.ndim == 2:
        image = cv2.cvtColor(image, cv2.COLOR_GRAY2BGRA)
    bgr = image[:, :, :3]
    height, width = bgr.shape[:2]
    if width < 3 or height < 3:
        raise ValueError("The image is too small for background removal.")
    mask = np.zeros((height, width), np.uint8)
    background_model = np.zeros((1, 65), np.float64)
    foreground_model = np.zeros((1, 65), np.float64)
    cv2.grabCut(bgr, mask, (1, 1, width - 2, height - 2), background_model, foreground_model, 5, cv2.GC_INIT_WITH_RECT)
    alpha = np.where((mask == cv2.GC_FGD) | (mask == cv2.GC_PR_FGD), 255, 0).astype(np.uint8)
    if image.shape[2] == 4:
        alpha = np.minimum(alpha, image[:, :, 3])
    result = np.dstack((bgr, alpha))
    partial = output_path.with_suffix(".partial.png")
    if not cv2.imwrite(str(partial), result):
        raise ValueError("Unable to save the transparent PNG.")
    os.replace(partial, output_path)


def capture_screenshot(
    output_path: Path,
    region: tuple[int, int, int, int] | None = None,
    annotation: str = "",
) -> Path:
    """Capture the full screen or a rectangular region, with an optional caption."""
    from PIL import ImageDraw, ImageFont, ImageGrab

    if output_path.suffix.casefold() not in {".png", ".jpg", ".jpeg", ".webp"}:
        raise ValueError("Screenshot output must be PNG, JPG, or WebP.")
    image = ImageGrab.grab(bbox=region, all_screens=True)
    try:
        if annotation.strip():
            draw = ImageDraw.Draw(image, "RGBA")
            font = ImageFont.load_default()
            left, top, right, bottom = draw.textbbox((0, 0), annotation.strip(), font=font)
            padding = 10
            draw.rounded_rectangle((10, 10, right - left + 20 + padding, bottom - top + 20 + padding), radius=6, fill=(0, 0, 0, 175))
            draw.text((20, 20), annotation.strip(), font=font, fill=(255, 255, 255, 255))
        partial = output_path.with_suffix(output_path.suffix + ".partial")
        image_format = {".png": "PNG", ".jpg": "JPEG", ".jpeg": "JPEG", ".webp": "WEBP"}[output_path.suffix.casefold()]
        if image_format == "JPEG" and image.mode not in {"RGB", "L"}:
            converted = image.convert("RGB")
            image.close()
            image = converted
        image.save(partial, format=image_format)
        os.replace(partial, output_path)
    finally:
        image.close()
    return output_path


def scan_document_from_webcam(output_path: Path, camera_index: int = 0) -> int:
    """Capture document pages from a webcam; Space captures, Enter finishes, Esc cancels."""
    import cv2
    from PIL import Image

    if output_path.suffix.casefold() != ".pdf":
        raise ValueError("Webcam scans must be saved as a PDF file.")
    camera = cv2.VideoCapture(camera_index)
    if not camera.isOpened():
        raise ValueError("No webcam was found. Connect a camera and try again.")
    window_name = "OpenMerger scanner — Space: capture | Enter: finish | Esc: cancel"
    captured: list[Path] = []
    try:
        with tempfile.TemporaryDirectory() as temp:
            while True:
                available, frame = camera.read()
                if not available:
                    raise ValueError("Unable to read from the webcam.")
                preview = _document_crop(frame)
                cv2.imshow(window_name, preview)
                key = cv2.waitKey(10) & 0xFF
                if key == 27:
                    raise ValueError("Webcam scan cancelled.")
                if key in {13, 10}:
                    break
                if key == 32:
                    target = Path(temp) / f"scan_{len(captured) + 1:04d}.jpg"
                    if not cv2.imwrite(str(target), preview):
                        raise ValueError("Unable to save the captured page.")
                    captured.append(target)
            if not captured:
                raise ValueError("Capture at least one page with the Space bar before finishing.")
            images = [Image.open(path).convert("RGB") for path in captured]
            try:
                partial = output_path.with_suffix(output_path.suffix + ".partial")
                images[0].save(partial, "PDF", save_all=True, append_images=images[1:], resolution=200)
                os.replace(partial, output_path)
            finally:
                for image in images:
                    image.close()
    finally:
        camera.release()
        cv2.destroyAllWindows()
    return len(captured)


def _document_crop(frame: object) -> object:
    """Return a perspective-corrected document when a four-corner page is visible."""
    import cv2

    height, width = frame.shape[:2]
    scale = 900 / max(height, width)
    small = cv2.resize(frame, (round(width * scale), round(height * scale)))
    gray = cv2.cvtColor(small, cv2.COLOR_BGR2GRAY)
    edges = cv2.Canny(cv2.GaussianBlur(gray, (5, 5), 0), 60, 180)
    contours, _ = cv2.findContours(edges, cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
    for contour in sorted(contours, key=cv2.contourArea, reverse=True):
        approximation = cv2.approxPolyDP(contour, 0.02 * cv2.arcLength(contour, True), True)
        if len(approximation) == 4 and cv2.contourArea(approximation) > small.shape[0] * small.shape[1] * 0.12:
            points = approximation.reshape(4, 2).astype("float32") / scale
            return _four_point_transform(frame, points)
    return frame


def _four_point_transform(image: object, points: object) -> object:
    import cv2
    import numpy as np

    ordered = np.zeros((4, 2), dtype="float32")
    sums = points.sum(axis=1)
    differences = np.diff(points, axis=1)
    ordered[0] = points[np.argmin(sums)]
    ordered[2] = points[np.argmax(sums)]
    ordered[1] = points[np.argmin(differences)]
    ordered[3] = points[np.argmax(differences)]
    top_left, top_right, bottom_right, bottom_left = ordered
    width = round(max(np.linalg.norm(bottom_right - bottom_left), np.linalg.norm(top_right - top_left)))
    height = round(max(np.linalg.norm(top_right - bottom_right), np.linalg.norm(top_left - bottom_left)))
    destination = np.array([[0, 0], [width - 1, 0], [width - 1, height - 1], [0, height - 1]], dtype="float32")
    return cv2.warpPerspective(image, cv2.getPerspectiveTransform(ordered, destination), (width, height))


def _register_heif_support() -> None:
    try:
        import pillow_heif
    except ImportError:
        return
    pillow_heif.register_heif_opener()


def add_page_numbers(source_path: Path, output_path: Path, position: str = "bottom-center", password: str | None = None) -> None:
    """Create a numbered PDF copy using a subtle footer on each page."""
    import pikepdf

    _ensure_new_output(source_path, output_path)
    with pikepdf.Pdf.open(source_path, password=password or "") as document:
        for number, page in enumerate(document.pages, 1):
            _add_text_overlay(page, f"Page {number} of {len(document.pages)}", position, 10, opacity=0.65)
        partial = output_path.with_suffix(output_path.suffix + ".partial")
        document.save(partial)
        partial.replace(output_path)


def add_text_watermark(source_path: Path, output_path: Path, text: str, password: str | None = None) -> None:
    """Add a centered translucent text watermark to every page."""
    import pikepdf

    if not text.strip():
        raise ValueError("Enter watermark text.")
    _ensure_new_output(source_path, output_path)
    with pikepdf.Pdf.open(source_path, password=password or "") as document:
        for page in document.pages:
            _add_text_overlay(page, text.strip(), "center", 28, opacity=0.22)
        partial = output_path.with_suffix(output_path.suffix + ".partial")
        document.save(partial)
        partial.replace(output_path)


def create_cover_page(output_path: Path, title: str, subtitle: str = "", page_size: tuple[int, int] = (595, 842)) -> None:
    """Create a clean, editable-text PDF cover page."""
    import pikepdf
    from pikepdf.canvas import BLUE, Canvas, Helvetica, Text

    if not title.strip():
        raise ValueError("Enter a cover-page title.")
    canvas = Canvas(page_size=page_size)
    canvas.add_font(pikepdf.Name("/F1"), Helvetica())
    width, height = page_size
    canvas.do.fill_color(BLUE)
    canvas.do.rect(0, height - 120, width, 120, fill=True)
    heading = Text().font(pikepdf.Name("/F1"), 30).move_cursor(54, height - 210).show(title.strip())
    canvas.do.draw_text(heading)
    if subtitle.strip():
        detail = Text().font(pikepdf.Name("/F1"), 15).move_cursor(56, height - 250).show(subtitle.strip())
        canvas.do.draw_text(detail)
    document = canvas.to_pdf()
    try:
        partial = output_path.with_suffix(output_path.suffix + ".partial")
        document.save(partial)
        partial.replace(output_path)
    finally:
        document.close()


def _add_text_overlay(page: object, text_value: str, position: str, size: int, opacity: float) -> None:
    import pikepdf
    from pikepdf.canvas import Canvas, Helvetica, Text

    media_box = page.mediabox
    width, height = float(media_box[2]), float(media_box[3])
    locations = {
        "bottom-center": (width * 0.42, 24),
        "bottom-right": (max(24, width - 110), 24),
        "top-right": (max(24, width - 110), height - 30),
        "center": (max(40, width * 0.28), height * 0.5),
    }
    x, y = locations.get(position, locations["bottom-center"])
    canvas = Canvas(page_size=(width, height))
    canvas.add_font(pikepdf.Name("/F1"), Helvetica())
    canvas.do.push()
    canvas.do.fill_color(pikepdf.canvas.Color(0.12, 0.23, 0.42, alpha=opacity))
    text = Text().font(pikepdf.Name("/F1"), size).move_cursor(x, y).show(text_value)
    canvas.do.draw_text(text)
    canvas.do.pop()
    overlay = canvas.to_pdf()
    try:
        page.add_overlay(overlay.pages[0])
    finally:
        overlay.close()


def pdf_to_images(source_path: Path, output_folder: Path, image_format: Literal["png", "jpg"] = "png", dpi: int = 150, password: str | None = None) -> list[Path]:
    """Render every PDF page to an image file."""
    import pymupdf

    if image_format not in {"png", "jpg"}:
        raise ValueError("Image format must be PNG or JPG.")
    if dpi not in {150, 300}:
        raise ValueError("DPI must be 150 or 300.")
    output_folder.mkdir(parents=True, exist_ok=True)
    document = pymupdf.open(source_path)
    if document.needs_pass:
        if not password or not document.authenticate(password):
            document.close()
            raise ValueError("A valid PDF password is required.")
    try:
        zoom = dpi / 72
        matrix = pymupdf.Matrix(zoom, zoom)
        output_paths: list[Path] = []
        for index, page in enumerate(document, 1):
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            target = output_folder / f"{source_path.stem}_page_{index:04d}.{image_format}"
            pixmap.save(target)
            output_paths.append(target)
        return output_paths
    finally:
        document.close()


def find_blank_pages(source_path: Path, threshold: float = 0.0001, password: str | None = None) -> list[int]:
    """Return 1-based page numbers whose low-resolution render is almost entirely white."""
    import pymupdf

    document = pymupdf.open(source_path)
    if document.needs_pass and (not password or not document.authenticate(password)):
        document.close()
        raise ValueError("A valid PDF password is required.")
    try:
        blanks: list[int] = []
        for index, page in enumerate(document, 1):
            pixmap = page.get_pixmap(matrix=pymupdf.Matrix(0.5, 0.5), colorspace=pymupdf.csGRAY, alpha=False)
            dark_pixels = sum(value < 245 for value in pixmap.samples)
            if dark_pixels / max(1, len(pixmap.samples)) <= threshold:
                blanks.append(index)
        return blanks
    finally:
        document.close()


def remove_blank_pages(source_path: Path, output_path: Path, threshold: float = 0.0001, password: str | None = None) -> list[int]:
    """Create a PDF copy with near-blank pages removed and return removed page numbers."""
    import pymupdf

    _ensure_new_output(source_path, output_path)
    blanks = find_blank_pages(source_path, threshold, password)
    document = pymupdf.open(source_path)
    if document.needs_pass:
        document.authenticate(password or "")
    try:
        if len(blanks) == document.page_count:
            raise ValueError("Every page appears blank; no output was created.")
        for page_number in reversed(blanks):
            document.delete_page(page_number - 1)
        document.save(output_path)
        return blanks
    finally:
        document.close()


def find_duplicate_pages(source_path: Path, password: str | None = None) -> list[list[int]]:
    """Find visually identical pages by hashing a consistent low-resolution rendering."""
    import pymupdf

    document = pymupdf.open(source_path)
    if document.needs_pass and (not password or not document.authenticate(password)):
        document.close()
        raise ValueError("A valid PDF password is required.")
    try:
        grouped: dict[str, list[int]] = {}
        for index, page in enumerate(document, 1):
            pixmap = page.get_pixmap(matrix=pymupdf.Matrix(0.5, 0.5), colorspace=pymupdf.csGRAY, alpha=False)
            grouped.setdefault(hashlib.sha256(pixmap.samples).hexdigest(), []).append(index)
        return [pages for pages in grouped.values() if len(pages) > 1]
    finally:
        document.close()


def pdf_to_docx(source_path: Path, output_path: Path, password: str | None = None) -> int:
    """Extract text into DOCX; scanned PDFs require OCR first."""
    import pymupdf
    from docx import Document

    _ensure_new_output(source_path, output_path)
    pdf = pymupdf.open(source_path)
    if pdf.needs_pass and (not password or not pdf.authenticate(password)):
        pdf.close()
        raise ValueError("A valid PDF password is required.")
    try:
        document = Document()
        count = 0
        for page in pdf:
            text = page.get_text().strip()
            if text:
                document.add_paragraph(text)
                count += 1
        if not count:
            raise ValueError("No selectable text found. Run OCR before converting this scanned PDF.")
        document.save(output_path)
        return count
    finally:
        pdf.close()


def pdf_to_pptx(source_path: Path, output_path: Path, dpi: int = 150, password: str | None = None) -> int:
    import pymupdf
    from pptx import Presentation

    _ensure_new_output(source_path, output_path)
    pdf = pymupdf.open(source_path)
    if pdf.needs_pass and (not password or not pdf.authenticate(password)):
        pdf.close()
        raise ValueError("A valid PDF password is required.")
    try:
        presentation = Presentation()
        presentation.slide_width = 13_333_200
        presentation.slide_height = 7_500_000
        with tempfile.TemporaryDirectory() as temp:
            for index, page in enumerate(pdf):
                pix = page.get_pixmap(matrix=pymupdf.Matrix(dpi / 72, dpi / 72), alpha=False)
                image = Path(temp) / f"page-{index}.png"
                pix.save(image)
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                slide.shapes.add_picture(str(image), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
        presentation.save(output_path)
        return pdf.page_count
    finally:
        pdf.close()


def excel_to_csv(source_path: Path, output_path: Path, sheet_name: str | None = None) -> int:
    from openpyxl import load_workbook

    _ensure_new_output(source_path, output_path)
    workbook = load_workbook(source_path, data_only=True, read_only=True)
    try:
        sheet = workbook[sheet_name] if sheet_name else workbook.active
        count = 0
        with output_path.open("w", newline="", encoding="utf-8-sig") as handle:
            writer = csv.writer(handle)
            for row in sheet.iter_rows(values_only=True):
                writer.writerow(row)
                count += 1
        return count
    finally:
        workbook.close()


def csv_to_excel(source_path: Path, output_path: Path) -> int:
    from openpyxl import Workbook

    _ensure_new_output(source_path, output_path)
    workbook = Workbook()
    sheet = workbook.active
    count = 0
    with source_path.open(newline="", encoding="utf-8-sig") as handle:
        for row in csv.reader(handle):
            sheet.append(row)
            count += 1
    workbook.save(output_path)
    return count


def find_office_executable(app_root: Path | None = None) -> Path | None:
    """Find the bundled conversion engine first, then a normal LibreOffice install."""
    roots: list[Path] = []
    if app_root is not None:
        roots.append(app_root)
    elif getattr(sys, "frozen", False):
        roots.append(Path(sys.executable).resolve().parent)
    else:
        roots.append(Path(__file__).resolve().parents[1])

    configured = os.environ.get("OPENMERGER_SOFFICE")
    candidates: list[Path] = [Path(configured)] if configured else []
    candidates.extend(root / "engine" / "LibreOffice" / "program" / "soffice.exe" for root in roots)
    for variable in ("ProgramFiles", "ProgramFiles(x86)"):
        directory = os.environ.get(variable)
        if directory:
            candidates.append(Path(directory) / "LibreOffice" / "program" / "soffice.exe")
    for candidate in candidates:
        if candidate.is_file():
            return candidate

    on_path = shutil.which("soffice") or shutil.which("libreoffice")
    return Path(on_path) if on_path else None


def office_to_pdf(source_path: Path, output_path: Path) -> None:
    """Convert DOCX/PPTX using OpenMerger's bundled or installed LibreOffice engine."""
    _ensure_new_output(source_path, output_path)
    office = find_office_executable()
    if not office:
        raise ValueError(
            "The document conversion engine is missing. Install the OpenMerger Full edition "
            "or install LibreOffice, then try again."
        )
    with tempfile.TemporaryDirectory() as temp:
        profile = (Path(temp) / "profile").resolve().as_uri()
        result = subprocess.run(
            [str(office), "--headless", f"-env:UserInstallation={profile}", "--convert-to", "pdf", "--outdir", temp, str(source_path)],
            capture_output=True,
            text=True,
            check=False,
            timeout=120,
        )
        created = Path(temp) / f"{source_path.stem}.pdf"
        if result.returncode or not created.exists():
            raise ValueError(result.stderr or "Office conversion failed.")
        shutil.move(created, output_path)


def parse_page_ranges(specification: str, page_count: int) -> list[int]:
    """Convert a 1-based range such as ``1-3,5,8-`` into zero-based page indexes."""
    indexes: list[int] = []
    for token in (part.strip() for part in specification.split(",")):
        if not token:
            continue
        if "-" in token:
            start_text, end_text = token.split("-", 1)
            start = int(start_text) if start_text else 1
            end = int(end_text) if end_text else page_count
            if start < 1 or end < start or end > page_count:
                raise ValueError(f"Invalid page range: {token}")
            indexes.extend(range(start - 1, end))
        else:
            number = int(token)
            if number < 1 or number > page_count:
                raise ValueError(f"Page {number} is outside this PDF")
            indexes.append(number - 1)
    if not indexes:
        raise ValueError("Select at least one page.")
    return indexes


def transform_pdf(source_path: Path, output_path: Path, pages: str | None = None, rotate: int = 0, password: str | None = None) -> int:
    """Extract/rotate pages into a new PDF without modifying the source."""
    import pikepdf

    if rotate % 90:
        raise ValueError("Rotation must be a multiple of 90 degrees.")
    _ensure_new_output(source_path, output_path)
    with pikepdf.Pdf.open(source_path, password=password or "") as source:
        selected = parse_page_ranges(pages, len(source.pages)) if pages else list(range(len(source.pages)))
        destination = pikepdf.Pdf.new()
        try:
            for index in selected:
                page = source.pages[index]
                destination.pages.append(page)
                if rotate:
                    copied = destination.pages[-1]
                    existing = int(copied.obj.get("/Rotate", 0))
                    copied.obj["/Rotate"] = (existing + rotate) % 360
            partial = output_path.with_suffix(output_path.suffix + ".partial")
            destination.save(partial)
            partial.replace(output_path)
        finally:
            destination.close()
    return len(selected)


def split_pdf(source_path: Path, output_path: Path, pages_per_file: int, password: str | None = None) -> list[Path]:
    import pikepdf

    if pages_per_file < 1:
        raise ValueError("Pages per file must be at least 1.")
    _ensure_new_output(source_path, output_path)
    with pikepdf.Pdf.open(source_path, password=password or "") as source:
        output_paths: list[Path] = []
        for start in range(0, len(source.pages), pages_per_file):
            target = output_path.with_name(f"{output_path.stem}_part_{start // pages_per_file + 1:04d}{output_path.suffix}")
            destination = pikepdf.Pdf.new()
            try:
                destination.pages.extend(source.pages[start : start + pages_per_file])
                partial = target.with_suffix(target.suffix + ".partial")
                destination.save(partial)
                partial.replace(target)
                output_paths.append(target)
            finally:
                destination.close()
    return output_paths


def images_to_pdf(
    image_paths: Iterable[Path],
    output_path: Path,
    page_preset: ImagePagePreset = "original",
    fit: ImageFit = "contain",
    margin_mm: float = 8,
    dpi: int = 150,
) -> int:
    """Create a print-ready PDF from images, with optional page-layout presets."""
    from PIL import Image

    if page_preset not in PAGE_SIZES_INCHES:
        raise ValueError("Choose a valid page preset.")
    if fit not in {"contain", "cover"}:
        raise ValueError("Image fit must be contain or cover.")
    if margin_mm < 0:
        raise ValueError("Margin cannot be negative.")
    if dpi not in {150, 300}:
        raise ValueError("DPI must be 150 or 300.")
    paths = list(image_paths)
    if output_path.expanduser().resolve() in {path.expanduser().resolve() for path in paths}:
        raise ValueError("Choose a different output file; source images are never overwritten.")
    images = []
    try:
        for path in paths:
            with Image.open(path) as source:
                image = source.convert("RGB")
                images.append(_layout_image(image, page_preset, fit, margin_mm, dpi))
                image.close()
        if not images:
            raise ValueError("Choose at least one image.")
        partial = output_path.with_suffix(output_path.suffix + ".partial")
        images[0].save(partial, "PDF", save_all=True, append_images=images[1:], resolution=float(dpi))
        partial.replace(output_path)
        return len(images)
    finally:
        for image in images:
            image.close()


def _layout_image(image: object, page_preset: ImagePagePreset, fit: ImageFit, margin_mm: float, dpi: int) -> object:
    from PIL import Image

    page_inches = PAGE_SIZES_INCHES[page_preset]
    if page_inches is None:
        return image.copy()
    page_width, page_height = page_inches
    if page_preset in {"a4", "letter", "poster"} and image.width > image.height and page_height > page_width:
        page_width, page_height = page_height, page_width
    canvas_size = (round(page_width * dpi), round(page_height * dpi))
    margin = round((margin_mm / 25.4) * dpi)
    available_width = max(1, canvas_size[0] - 2 * margin)
    available_height = max(1, canvas_size[1] - 2 * margin)
    scale = max(available_width / image.width, available_height / image.height) if fit == "cover" else min(available_width / image.width, available_height / image.height)
    resized = image.resize((max(1, round(image.width * scale)), max(1, round(image.height * scale))), Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", canvas_size, "white")
    x = margin + (available_width - resized.width) // 2
    y = margin + (available_height - resized.height) // 2
    canvas.paste(resized, (x, y))
    resized.close()
    return canvas


def update_metadata(source_path: Path, output_path: Path, title: str = "", author: str = "", subject: str = "", password: str | None = None) -> None:
    import pikepdf

    _ensure_new_output(source_path, output_path)
    with pikepdf.Pdf.open(source_path, password=password or "") as document:
        if title:
            document.docinfo["/Title"] = title
        if author:
            document.docinfo["/Author"] = author
        if subject:
            document.docinfo["/Subject"] = subject
        partial = output_path.with_suffix(output_path.suffix + ".partial")
        document.save(partial)
        partial.replace(output_path)


def protect_pdf(source_path: Path, output_path: Path, user_password: str, owner_password: str | None = None, source_password: str | None = None) -> None:
    import pikepdf

    if not user_password:
        raise ValueError("Enter a password to protect the PDF.")
    _ensure_new_output(source_path, output_path)
    with pikepdf.Pdf.open(source_path, password=source_password or "") as document:
        partial = output_path.with_suffix(output_path.suffix + ".partial")
        document.save(partial, encryption=pikepdf.Encryption(user=user_password, owner=owner_password or user_password, R=6))
        partial.replace(output_path)


def unlock_pdf(source_path: Path, output_path: Path, password: str) -> None:
    import pikepdf

    if not password:
        raise ValueError("Enter the current PDF password to unlock it.")
    _ensure_new_output(source_path, output_path)
    with pikepdf.Pdf.open(source_path, password=password) as document:
        partial = output_path.with_suffix(output_path.suffix + ".partial")
        document.save(partial, encryption=False)
        partial.replace(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="Non-destructive OpenMerger PDF page tools")
    parser.add_argument("source", type=Path)
    parser.add_argument("output", type=Path)
    parser.add_argument("--pages", help="1-based ranges, e.g. 1-3,5,8-")
    parser.add_argument("--rotate", type=int, default=0, help="Rotation in degrees (multiple of 90)")
    parser.add_argument("--split", type=int, metavar="PAGES_PER_FILE")
    parser.add_argument("--password", help="Password for an encrypted input; never saved")
    args = parser.parse_args()
    if args.split:
        outputs = split_pdf(args.source, args.output, args.split, args.password)
        print("\n".join(str(path) for path in outputs))
    else:
        count = transform_pdf(args.source, args.output, args.pages, args.rotate, args.password)
        print(f"Created {args.output} with {count} page(s).")


if __name__ == "__main__":
    main()
